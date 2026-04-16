"""PPTX export service for classroom-facing presentation projects."""

from __future__ import annotations

from dataclasses import dataclass
from io import BytesIO
from pathlib import Path
from typing import Any
from urllib.parse import unquote, urlparse
import re

from ..core.settings import settings
from ..models import Plan
from ..presentation_layouts import (
    BoxSpec,
    PresentationRenderContext,
    get_presentation_renderer,
    get_presentation_template,
    register_presentation_renderer,
)
from ..presentation_models import PresentationDocument, Slide, normalize_slide_template, paginate_slide_text
from ..presentation_style import (
    extract_presentation_style,
    get_theme_palette,
    normalize_presentation_style,
    resolve_density_limits,
    resolve_font_sizes,
)
from .presentation_service import PresentationService


class PresentationExportError(Exception):
    """Base error for presentation export failures."""


class PresentationNotFoundError(PresentationExportError):
    """Raised when the requested presentation project does not exist."""


class PresentationExportUnavailableError(PresentationExportError):
    """Raised when python-pptx is unavailable."""


SLIDE_WIDTH_INCHES = 10
SLIDE_HEIGHT_INCHES = 7.5
PROJECT_ROOT = Path(__file__).resolve().parents[3]
ABSOLUTE_URL_PATTERN = re.compile(r"(https?://[^\s]+)")
UPLOADS_URL_PATTERN = re.compile(r"((?:/?uploads/games/[^\s]+))")


@dataclass(frozen=True)
class RenderStyle:
    """Resolved style data shared across all rendered slides."""

    palette: Any
    school_name: str | None
    logo_path: Path | None
    density: str


def export_to_pptx(
    presentation_data: PresentationDocument | dict[str, Any],
    presentation_style: Any | None = None,
) -> bytes:
    """Render a presentation document into a `.pptx` binary."""
    try:
        from pptx import Presentation as PPTXPresentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:  # pragma: no cover - depends on local environment.
        raise PresentationExportUnavailableError(
            "PPTX 导出依赖未安装，请在 lesson_agent_env 中安装 python-pptx。"
        ) from exc

    document = (
        presentation_data
        if isinstance(presentation_data, PresentationDocument)
        else PresentationDocument.model_validate(presentation_data)
    )
    normalized_style = normalize_presentation_style(presentation_style)
    render_style = RenderStyle(
        palette=get_theme_palette(normalized_style.theme),
        school_name=(normalized_style.school_name or "").strip() or None,
        logo_path=_resolve_image_path(normalized_style.logo_url),
        density=normalized_style.density,
    )

    prs = PPTXPresentation()
    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)

    blank_layout_index = 6 if len(prs.slide_layouts) > 6 else max(len(prs.slide_layouts) - 1, 0)
    blank_layout = prs.slide_layouts[blank_layout_index]

    rendered_pages = (
        [Slide(title=document.title, body="", template="title_body")]
        if not document.slides
        else [page for slide_data in document.slides for page in _paginate_slide(slide_data, render_style=render_style)]
    )

    for page_number, slide_data in enumerate(rendered_pages, start=1):
        slide = prs.slides.add_slide(blank_layout)
        _render_slide(
            slide=slide,
            slide_data=slide_data,
            render_style=render_style,
            page_number=page_number,
            total_slides=len(rendered_pages),
            rgb=RGBColor,
            shape_type=MSO_AUTO_SHAPE_TYPE,
            inches=Inches,
            pt=Pt,
            align=PP_ALIGN,
            anchor=MSO_ANCHOR,
        )

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


class PresentationExportService:
    """Generate PPTX downloads for presentation projects."""

    PPTX_MEDIA_TYPE = "application/vnd.openxmlformats-officedocument.presentationml.presentation"

    def __init__(self, presentation_service: PresentationService) -> None:
        self.presentation_service = presentation_service

    def export_to_pptx(self, plan_id: str) -> bytes:
        """Load a presentation project by id and export it."""
        plan = self.presentation_service.get(plan_id)
        if plan is None:
            raise PresentationNotFoundError("演示文稿不存在。")
        return self.render_plan_to_pptx(plan)

    def render_plan_to_pptx(self, plan: Plan) -> bytes:
        """Render an already loaded presentation project."""
        content = plan.content if isinstance(plan.content, dict) else {"slides": []}
        metadata = getattr(plan, "metadata_json", None)
        if metadata is None:
            metadata = getattr(plan, "metadata", None)
        payload = {
            "title": content.get("title") or plan.title,
            "classroom_script": content.get("classroom_script", ""),
            "slides": content.get("slides", []),
        }
        return export_to_pptx(payload, presentation_style=extract_presentation_style(metadata))


def _paginate_slide(slide_data: Slide, *, render_style: RenderStyle) -> list[Slide]:
    """Split a long slide body into multiple pages while keeping the title stable."""
    template = normalize_slide_template(slide_data.template, slide_data.layout)
    template_spec = get_presentation_template(template)
    normalized_body, link_text, link_url = _normalize_slide_body_link_fields(
        body=slide_data.body,
        link_text=slide_data.link_text,
        link_url=slide_data.link_url,
    )
    if template_spec.renderer_name == "title_subtitle":
        return [
            slide_data.model_copy(
                update={
                    "template": template,
                    "layout": template_spec.pptx_layout,
                    "body": "",
                    "bullet_points": [],
                }
            )
        ]

    chars_per_line, max_lines = resolve_density_limits(
        density=render_style.density,
        chars_per_line=template_spec.pagination.chars_per_line,
        max_lines=template_spec.pagination.max_lines,
        has_image_panel=template_spec.image_box is not None,
    )
    pages = _paginate_text(normalized_body, line_width=chars_per_line, max_lines=max_lines)
    return [
        slide_data.model_copy(
            update={
                "template": template,
                "layout": template_spec.pptx_layout,
                "title": f"{slide_data.title}（{index + 1}/{len(pages)}）" if len(pages) > 1 else slide_data.title,
                "body": page,
                "bullet_points": [line for line in page.splitlines() if line.strip()],
                "link_text": link_text,
                "link_url": link_url,
            }
        )
        for index, page in enumerate(pages)
    ] or [slide_data]


def _paginate_text(text: str, *, line_width: int, max_lines: int) -> list[str]:
    """Wrap Chinese text approximately by characters and split into multiple pages."""
    return paginate_slide_text(text, chars_per_line=line_width, max_lines=max_lines)


def _render_slide(
    *,
    slide: Any,
    slide_data: Slide,
    render_style: RenderStyle,
    page_number: int,
    total_slides: int,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
) -> None:
    """Render one slide through the registered layout renderer."""
    template = normalize_slide_template(slide_data.template, slide_data.layout)
    template_spec = get_presentation_template(template)
    fonts = resolve_font_sizes(density=render_style.density, has_image_panel=template_spec.image_box is not None)

    if template_spec.renderer_name == "title_subtitle":
        _paint_title_subtitle_background(
            slide=slide,
            render_style=render_style,
            rgb=rgb,
            shape_type=shape_type,
            inches=inches,
            pt=pt,
            align=align,
            anchor=anchor,
        )
    else:
        _paint_background(
            slide=slide,
            render_style=render_style,
            page_number=page_number,
            total_slides=total_slides,
            rgb=rgb,
            shape_type=shape_type,
            inches=inches,
            pt=pt,
            align=align,
            anchor=anchor,
            branding_font_size=fonts["branding"],
        )
        _add_title_box(
            slide=slide,
            title=slide_data.title,
            render_style=render_style,
            font_size=fonts["title"],
            rgb=rgb,
            inches=inches,
            pt=pt,
            align=align,
            anchor=anchor,
        )

    renderer = get_presentation_renderer(template_spec.renderer_name)
    renderer(
        PresentationRenderContext(
            slide=slide,
            slide_data=slide_data,
            template_spec=template_spec,
            rgb=rgb,
            shape_type=shape_type,
            inches=inches,
            pt=pt,
            align=align,
            anchor=anchor,
            render_options={"style": render_style, "fonts": fonts},
        )
    )
    _set_notes(slide, slide_data.notes)


def _render_body_template(context: PresentationRenderContext) -> None:
    """Render a title + body layout."""
    render_style: RenderStyle = context.render_options["style"]
    fonts: dict[str, int] = context.render_options["fonts"]
    _add_body_box(
        slide=context.slide,
        text=context.slide_data.body,
        link_text=context.slide_data.link_text,
        link_url=context.slide_data.link_url,
        render_style=render_style,
        rgb=context.rgb,
        shape_type=context.shape_type,
        inches=context.inches,
        pt=context.pt,
        align=context.align,
        anchor=context.anchor,
        box=context.template_spec.body_box,
        font_size=fonts["body"],
    )


def _render_body_image_template(context: PresentationRenderContext) -> None:
    """Render a title + body + image layout."""
    _render_body_template(context)
    if context.template_spec.image_box is None:
        return
    render_style: RenderStyle = context.render_options["style"]
    fonts: dict[str, int] = context.render_options["fonts"]
    _add_image_panel(
        slide=context.slide,
        image_url=context.slide_data.image_url,
        image_description=context.slide_data.image_description,
        render_style=render_style,
        rgb=context.rgb,
        shape_type=context.shape_type,
        inches=context.inches,
        pt=context.pt,
        align=context.align,
        anchor=context.anchor,
        box=context.template_spec.image_box,
        placeholder_font_size=fonts["placeholder"],
        branding_font_size=fonts["branding"],
    )


def _render_title_subtitle_template(context: PresentationRenderContext) -> None:
    """Render a centered title slide with an optional subtitle."""
    render_style: RenderStyle = context.render_options["style"]
    fonts: dict[str, int] = context.render_options["fonts"]
    _add_centered_title_block(
        slide=context.slide,
        title=context.slide_data.title,
        subtitle=_resolve_slide_subtitle(context.slide_data),
        render_style=render_style,
        rgb=context.rgb,
        shape_type=context.shape_type,
        inches=context.inches,
        pt=context.pt,
        align=context.align,
        anchor=context.anchor,
        subtitle_box=context.template_spec.body_box,
        title_font_size=fonts["cover_title"],
        subtitle_font_size=fonts["subtitle"],
    )


def _paint_background(
    *,
    slide: Any,
    render_style: RenderStyle,
    page_number: int,
    total_slides: int,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
    branding_font_size: int,
) -> None:
    """Draw the shared classroom slide background and brand header."""
    palette = render_style.palette
    bg = slide.shapes.add_shape(shape_type.RECTANGLE, 0, 0, inches(SLIDE_WIDTH_INCHES), inches(SLIDE_HEIGHT_INCHES))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(*palette.background)
    bg.line.fill.background()

    header = slide.shapes.add_shape(shape_type.RECTANGLE, 0, 0, inches(SLIDE_WIDTH_INCHES), inches(1.45))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb(*palette.header)
    header.line.fill.background()

    accent = slide.shapes.add_shape(shape_type.RECTANGLE, 0, inches(1.45), inches(SLIDE_WIDTH_INCHES), inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(*palette.accent)
    accent.line.fill.background()

    if render_style.school_name:
        _add_text_box(
            slide=slide,
            text=render_style.school_name,
            left=6.1,
            top=0.44,
            width=2.55,
            height=0.32,
            font_size=branding_font_size,
            text_color=palette.title_on_header,
            rgb=rgb,
            shape_type=shape_type,
            inches=inches,
            pt=pt,
            anchor=anchor,
            paragraph_align=align.RIGHT,
        )

    if render_style.logo_path and render_style.logo_path.exists():
        _add_fitted_image(slide, render_style.logo_path, inches, BoxSpec(left=8.86, top=0.28, width=0.58, height=0.58))

    _add_text_box(
        slide=slide,
        text=f"{page_number}/{total_slides}",
        left=9.04,
        top=7.03,
        width=0.48,
        height=0.22,
        font_size=branding_font_size,
        text_color=palette.subtitle,
        rgb=rgb,
        shape_type=shape_type,
        inches=inches,
        pt=pt,
        anchor=anchor,
        paragraph_align=align.RIGHT,
    )


def _paint_title_subtitle_background(
    *,
    slide: Any,
    render_style: RenderStyle,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
) -> None:
    """Draw a calmer full-page layout for cover and ending slides."""
    palette = render_style.palette
    bg = slide.shapes.add_shape(shape_type.RECTANGLE, 0, 0, inches(SLIDE_WIDTH_INCHES), inches(SLIDE_HEIGHT_INCHES))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb(*palette.cover_background)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(shape_type.RECTANGLE, 0, 0, inches(1.15), inches(SLIDE_HEIGHT_INCHES))
    accent.fill.solid()
    accent.fill.fore_color.rgb = rgb(*palette.header)
    accent.line.fill.background()

    footer = slide.shapes.add_shape(shape_type.RECTANGLE, inches(0.95), inches(6.45), inches(8.15), inches(0.16))
    footer.fill.solid()
    footer.fill.fore_color.rgb = rgb(*palette.accent)
    footer.line.fill.background()

    if render_style.school_name:
        _add_text_box(
            slide=slide,
            text=render_style.school_name,
            left=1.28,
            top=5.9,
            width=5.8,
            height=0.35,
            font_size=resolve_font_sizes(density=render_style.density, has_image_panel=False)["branding"],
            text_color=palette.subtitle,
            rgb=rgb,
            shape_type=shape_type,
            inches=inches,
            pt=pt,
            anchor=anchor,
            paragraph_align=align.LEFT,
        )


def _add_title_box(
    *,
    slide: Any,
    title: str,
    render_style: RenderStyle,
    font_size: int,
    rgb: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
) -> None:
    """Render the top title with a stable two-line limit."""
    textbox = slide.shapes.add_textbox(inches(0.8), inches(0.42), inches(7.95), inches(0.82))
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = _truncate_title(title)
    paragraph.alignment = align.LEFT
    run = _get_or_add_first_run(paragraph)
    run.font.name = "Microsoft YaHei"
    run.font.size = pt(font_size)
    run.font.bold = True
    run.font.color.rgb = rgb(*render_style.palette.title_on_header)


def _truncate_title(title: str) -> str:
    normalized = (title or "").strip() or "未命名页面"
    if len(normalized) <= 28:
        return normalized
    return f"{normalized[:27]}…"


def _resolve_slide_subtitle(slide_data: Slide) -> str:
    """Prefer the explicit subtitle, then fall back to body for legacy title slides."""
    subtitle = str(slide_data.subtitle or "").strip()
    if subtitle:
        return subtitle
    return str(slide_data.body or "").strip()


def _add_centered_title_block(
    *,
    slide: Any,
    title: str,
    subtitle: str,
    render_style: RenderStyle,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
    subtitle_box: BoxSpec,
    title_font_size: int,
    subtitle_font_size: int,
) -> None:
    """Render the centered title and optional subtitle for cover/end slides."""
    title_box = slide.shapes.add_textbox(inches(1.25), inches(2.05), inches(7.35), inches(1.65))
    title_frame = title_box.text_frame
    title_frame.clear()
    title_frame.word_wrap = True
    title_frame.vertical_anchor = anchor.MIDDLE
    title_paragraph = title_frame.paragraphs[0]
    title_paragraph.text = (title or "").strip() or "未命名页面"
    title_paragraph.alignment = align.CENTER
    title_run = _get_or_add_first_run(title_paragraph)
    title_run.font.name = "Microsoft YaHei"
    title_run.font.size = pt(title_font_size)
    title_run.font.bold = True
    title_run.font.color.rgb = rgb(*render_style.palette.title_on_cover)

    if not subtitle:
        return

    subtitle_panel = slide.shapes.add_shape(
        shape_type.ROUNDED_RECTANGLE,
        inches(subtitle_box.left),
        inches(subtitle_box.top),
        inches(subtitle_box.width),
        inches(subtitle_box.height),
    )
    subtitle_panel.fill.solid()
    subtitle_panel.fill.fore_color.rgb = rgb(*render_style.palette.surface)
    subtitle_panel.line.color.rgb = rgb(*render_style.palette.border)
    subtitle_frame = subtitle_panel.text_frame
    subtitle_frame.clear()
    subtitle_frame.word_wrap = True
    subtitle_frame.vertical_anchor = anchor.MIDDLE
    subtitle_frame.margin_left = inches(0.16)
    subtitle_frame.margin_right = inches(0.16)
    subtitle_paragraph = subtitle_frame.paragraphs[0]
    subtitle_paragraph.text = subtitle
    subtitle_paragraph.alignment = align.CENTER
    subtitle_run = _get_or_add_first_run(subtitle_paragraph)
    subtitle_run.font.name = "Microsoft YaHei"
    subtitle_run.font.size = pt(subtitle_font_size)
    subtitle_run.font.color.rgb = rgb(*render_style.palette.subtitle)


def _add_body_box(
    *,
    slide: Any,
    text: str,
    link_text: str | None,
    link_url: str | None,
    render_style: RenderStyle,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
    box: BoxSpec,
    font_size: int,
) -> None:
    """Render the main classroom text area."""
    panel = slide.shapes.add_shape(
        shape_type.ROUNDED_RECTANGLE,
        inches(box.left),
        inches(box.top),
        inches(box.width),
        inches(box.height),
    )
    panel.fill.solid()
    panel.fill.fore_color.rgb = rgb(*render_style.palette.surface)
    panel.line.color.rgb = rgb(*render_style.palette.border)
    frame = panel.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor.TOP
    frame.margin_left = inches(0.18)
    frame.margin_right = inches(0.18)
    frame.margin_top = inches(0.18)
    frame.margin_bottom = inches(0.18)

    lines = text.splitlines() if text.strip() else ["（本页暂无正文）"]
    first = True
    for line in lines:
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        first = False
        paragraph.text = ""
        paragraph.alignment = align.LEFT
        paragraph.space_after = pt(6)
        _populate_body_paragraph_with_links(
            paragraph=paragraph,
            text=line,
            link_text=link_text,
            link_url=link_url,
            render_style=render_style,
            rgb=rgb,
            pt=pt,
            font_size=font_size,
        )


def _add_image_panel(
    *,
    slide: Any,
    image_url: str | None,
    image_description: str | None,
    render_style: RenderStyle,
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    align: Any,
    anchor: Any,
    box: BoxSpec,
    placeholder_font_size: int,
    branding_font_size: int,
) -> None:
    """Render either the uploaded image or a labeled placeholder box."""
    container = slide.shapes.add_shape(
        shape_type.ROUNDED_RECTANGLE,
        inches(box.left),
        inches(box.top),
        inches(box.width),
        inches(box.height),
    )
    container.fill.solid()
    container.fill.fore_color.rgb = rgb(*render_style.palette.surface)
    container.line.color.rgb = rgb(*render_style.palette.border)

    image_path = _resolve_image_path(image_url)
    if image_path and image_path.exists():
        _add_fitted_image(slide, image_path, inches, box)
        return

    textbox = slide.shapes.add_textbox(
        inches(box.left + 0.18),
        inches(box.top + 0.2),
        inches(max(box.width - 0.36, 0.2)),
        inches(max(box.height - 0.4, 0.2)),
    )
    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor.MIDDLE
    paragraph = frame.paragraphs[0]
    placeholder_label = "图片占位"
    description = (image_description or "").strip()
    paragraph.text = placeholder_label if not description else f"{placeholder_label}\n{description}"
    paragraph.alignment = align.CENTER
    run = _get_or_add_first_run(paragraph)
    run.font.name = "Microsoft YaHei"
    run.font.size = pt(placeholder_font_size)
    run.font.color.rgb = rgb(*render_style.palette.subtitle)


def _add_text_box(
    *,
    slide: Any,
    text: str,
    left: float,
    top: float,
    width: float,
    height: float,
    font_size: int,
    text_color: tuple[int, int, int],
    rgb: Any,
    shape_type: Any,
    inches: Any,
    pt: Any,
    anchor: Any,
    paragraph_align: Any,
    fill_color: tuple[int, int, int] | None = None,
    fill_transparency: float = 0.0,
) -> None:
    """Add a compact rounded text box used by branding and captions."""
    textbox = slide.shapes.add_shape(
        shape_type.ROUNDED_RECTANGLE,
        inches(left),
        inches(top),
        inches(width),
        inches(height),
    )
    if fill_color is None:
        textbox.fill.background()
    else:
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = rgb(*fill_color)
        textbox.fill.transparency = fill_transparency
    textbox.line.fill.background()

    frame = textbox.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = anchor.MIDDLE
    frame.margin_left = inches(0.08)
    frame.margin_right = inches(0.08)
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = paragraph_align
    run = _get_or_add_first_run(paragraph)
    run.font.name = "Microsoft YaHei"
    run.font.size = pt(font_size)
    run.font.color.rgb = rgb(*text_color)


def _get_or_add_first_run(paragraph: Any) -> Any:
    """Return the first text run, creating one for empty paragraphs when needed."""
    runs = getattr(paragraph, "runs", ())
    if runs:
        return runs[0]
    return paragraph.add_run()


def _populate_body_paragraph_with_links(
    *,
    paragraph: Any,
    text: str,
    link_text: str | None,
    link_url: str | None,
    render_style: RenderStyle,
    rgb: Any,
    pt: Any,
    font_size: int,
) -> None:
    """Render one body line and attach hyperlinks for absolute URLs when present."""
    if link_url and link_text and text.strip() == link_text.strip():
        run = paragraph.add_run()
        run.text = text
        run.font.name = "Microsoft YaHei"
        run.font.size = pt(font_size)
        run.font.color.rgb = rgb(5, 99, 193)
        run.font.underline = True
        run.hyperlink.address = link_url
        return

    segments = _split_text_segments(text)
    if not segments:
        segments = [(text, None)]

    for segment_text, hyperlink in segments:
        run = paragraph.add_run()
        run.text = segment_text
        run.font.name = "Microsoft YaHei"
        run.font.size = pt(font_size)
        if hyperlink:
            run.font.color.rgb = rgb(5, 99, 193)
            run.font.underline = True
            run.hyperlink.address = hyperlink
        else:
            run.font.color.rgb = rgb(*render_style.palette.body)


def _split_text_segments(text: str) -> list[tuple[str, str | None]]:
    """Split one line into plain-text and hyperlink segments."""
    if not text:
        return []

    segments: list[tuple[str, str | None]] = []
    cursor = 0
    for match in ABSOLUTE_URL_PATTERN.finditer(text):
        start, end = match.span()
        if start > cursor:
            segments.append((text[cursor:start], None))
        url = match.group(1)
        segments.append((url, url))
        cursor = end
    if cursor < len(text):
        segments.append((text[cursor:], None))
    return segments


def _normalize_slide_body_link_fields(
    *,
    body: str,
    link_text: str | None,
    link_url: str | None,
) -> tuple[str, str | None, str | None]:
    """Keep long links out of paginated body text while preserving a clickable entry line."""
    resolved_link_url = _resolve_public_link_url(link_url)
    resolved_link_text = str(link_text or "").strip() or None
    normalized_lines: list[str] = []

    for raw_line in body.splitlines():
        line = str(raw_line or "")
        candidate_url = _extract_first_line_url(line)
        if not candidate_url:
            normalized_lines.append(line)
            continue

        normalized_url = _resolve_public_link_url(candidate_url)
        if not resolved_link_url and normalized_url:
            resolved_link_url = normalized_url

        replacement_text = resolved_link_text or _build_link_line_text(line, candidate_url)
        resolved_link_text = replacement_text
        normalized_lines.append(replacement_text)

    if resolved_link_url and resolved_link_text:
        normalized_text = "\n".join(normalized_lines).strip()
        if resolved_link_text not in normalized_text.splitlines():
            normalized_lines.append(resolved_link_text)

    normalized_body = "\n".join(normalized_lines).strip()
    return normalized_body, resolved_link_text, resolved_link_url


def _extract_first_line_url(text: str) -> str | None:
    """Extract one URL-like target from a body line."""
    for pattern in (ABSOLUTE_URL_PATTERN, UPLOADS_URL_PATTERN):
        match = pattern.search(text)
        if match:
            return _strip_trailing_url_punctuation(match.group(1))
    return None


def _build_link_line_text(text: str, url: str) -> str:
    """Replace a raw URL with a compact classroom-facing label."""
    stripped = text.strip()
    for separator in ("：", ":"):
        if separator in stripped:
            prefix, _ = stripped.split(separator, 1)
            if prefix.strip() in {"互动入口", "互动页面"}:
                return f"{prefix.strip()}：点击打开小游戏"
    if stripped == url:
        return "点击打开链接"
    replaced = stripped.replace(url, "点击打开链接").strip()
    return replaced or "点击打开链接"


def _resolve_public_link_url(value: str | None) -> str | None:
    """Upgrade stored relative link targets to absolute public URLs."""
    raw_value = str(value or "").strip()
    if not raw_value:
        return None
    raw_value = _strip_trailing_url_punctuation(raw_value)
    if raw_value.startswith(("http://", "https://")):
        return raw_value
    if raw_value.startswith("/"):
        return f"{settings.PUBLIC_BASE_URL}{raw_value}"
    if raw_value.startswith("uploads/"):
        return f"{settings.PUBLIC_BASE_URL}/{raw_value}"
    return raw_value


def _strip_trailing_url_punctuation(value: str) -> str:
    """Trim sentence punctuation accidentally captured alongside a URL."""
    return value.rstrip("。，、；；：:!?！？）)]】>,")


def _resolve_image_path(image_url: str | None) -> Path | None:
    """Map stored image URLs back to a readable local file path when possible."""
    raw_value = str(image_url or "").strip()
    if not raw_value:
        return None

    direct_path = Path(raw_value).expanduser()
    if direct_path.exists():
        return direct_path

    parsed = urlparse(raw_value)
    decoded_path = unquote((parsed.path or raw_value).strip())
    if not decoded_path:
        return None

    if parsed.scheme == "file":
        candidate = Path(decoded_path).expanduser()
        return candidate if candidate.exists() else None

    if decoded_path.startswith("/uploads/"):
        candidate = PROJECT_ROOT / decoded_path.lstrip("/")
        return candidate if candidate.exists() else None

    if decoded_path.startswith("uploads/"):
        candidate = PROJECT_ROOT / decoded_path
        return candidate if candidate.exists() else None

    candidate = Path(decoded_path).expanduser()
    return candidate if candidate.exists() else None


def _add_fitted_image(slide: Any, image_path: Path, inches: Any, box: BoxSpec) -> None:
    """Insert an image while preserving aspect ratio inside the target panel."""
    try:
        from PIL import Image
    except ImportError:
        slide.shapes.add_picture(str(image_path), inches(box.left), inches(box.top), width=inches(box.width))
        return

    with Image.open(image_path) as image:
        img_width, img_height = image.size

    if img_width <= 0 or img_height <= 0:
        slide.shapes.add_picture(str(image_path), inches(box.left), inches(box.top), width=inches(box.width))
        return

    box_ratio = box.width / box.height
    image_ratio = img_width / img_height
    if image_ratio >= box_ratio:
        final_width = box.width
        final_height = box.width / image_ratio
        offset_x = box.left
        offset_y = box.top + (box.height - final_height) / 2
    else:
        final_height = box.height
        final_width = box.height * image_ratio
        offset_x = box.left + (box.width - final_width) / 2
        offset_y = box.top

    slide.shapes.add_picture(
        str(image_path),
        inches(offset_x),
        inches(offset_y),
        width=inches(final_width),
        height=inches(final_height),
    )


def _set_notes(slide: Any, notes: str | None) -> None:
    """Write speaker notes when provided."""
    if not notes or not notes.strip():
        return

    try:
        slide.notes_slide.notes_text_frame.text = notes.strip()
    except Exception:  # noqa: BLE001
        return


register_presentation_renderer("body", _render_body_template)
register_presentation_renderer("body_image", _render_body_image_template)
register_presentation_renderer("title_subtitle", _render_title_subtitle_template)
