"""Regression tests for PPTX export edge cases."""

from __future__ import annotations

import tempfile
import unittest
import uuid
from pathlib import Path
import sys

PROJECT_ROOT = Path(__file__).resolve().parents[2]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from backend.app.services.export_pptx import export_to_pptx

UPLOADS_DIR = PROJECT_ROOT / "uploads" / "images"
TINY_PNG = (
    b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
    b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDAT\x08\xd7c\xf8\xff\xff?"
    b"\x00\x05\xfe\x02\xfeA\x8d\x89\xb5\x00\x00\x00\x00IEND\xaeB`\x82"
)


class ExportPptxTests(unittest.TestCase):
    """Protect image embedding behavior used by the PPT editor/export flow."""

    def test_export_allows_blank_lines_in_slide_body(self) -> None:
        pptx_bytes = export_to_pptx(
            {
                "title": "空行导出",
                "classroom_script": "",
                "slides": [
                    {
                        "template": "title_body",
                        "title": "课堂流程",
                        "body": "第一步：观察现象。\n\n第二步：提出猜想。",
                    }
                ],
            }
        )

        self.assertGreater(len(pptx_bytes), 0)
        self.assertEqual(pptx_bytes[:2], b"PK")

    def test_export_writes_clickable_hyperlink_for_absolute_game_url(self) -> None:
        pptx_bytes = export_to_pptx(
            {
                "title": "小游戏链接",
                "classroom_script": "",
                "slides": [
                    {
                        "template": "title_body",
                        "title": "浮力快答",
                        "body": (
                            "根据课堂内容完成选择题。\n"
                            "互动入口：http://127.0.0.1:8000/uploads/games/game_demo.html"
                        ),
                    }
                ],
            }
        )

        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            temp_file.write(pptx_bytes)
            temp_file.flush()

            from pptx import Presentation as PPTXPresentation

            exported = PPTXPresentation(temp_file.name)
            hyperlink_targets = [
                run.hyperlink.address
                for shape in exported.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if getattr(run, "hyperlink", None) and run.hyperlink.address
            ]

            self.assertIn("http://127.0.0.1:8000/uploads/games/game_demo.html", hyperlink_targets)

    def test_export_writes_clickable_hyperlink_from_dedicated_link_field(self) -> None:
        pptx_bytes = export_to_pptx(
            {
                "title": "小游戏链接",
                "classroom_script": "",
                "slides": [
                    {
                        "template": "title_body",
                        "title": "浮力快答",
                        "body": "根据课堂内容完成选择题。\n巩固目标：判断浮力变化\n互动入口：点击打开小游戏",
                        "link_text": "互动入口：点击打开小游戏",
                        "link_url": "http://127.0.0.1:8000/uploads/games/game_demo.html",
                    }
                ],
            }
        )

        with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
            temp_file.write(pptx_bytes)
            temp_file.flush()

            from pptx import Presentation as PPTXPresentation

            exported = PPTXPresentation(temp_file.name)
            hyperlink_runs = [
                run
                for shape in exported.slides[0].shapes
                if getattr(shape, "has_text_frame", False)
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if getattr(run, "hyperlink", None) and run.hyperlink.address
            ]

            self.assertTrue(any(run.text == "互动入口：点击打开小游戏" for run in hyperlink_runs))
            self.assertTrue(
                any(
                    run.hyperlink.address == "http://127.0.0.1:8000/uploads/games/game_demo.html"
                    for run in hyperlink_runs
                )
            )

    def test_export_embeds_image_for_public_uploads_url(self) -> None:
        UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"ppt-export-{uuid.uuid4().hex}.png"
        image_path = UPLOADS_DIR / filename
        image_path.write_bytes(TINY_PNG)

        try:
            pptx_bytes = export_to_pptx(
                {
                    "title": "带图导出",
                    "classroom_script": "",
                    "slides": [
                        {
                            "template": "title_body_image",
                            "title": "实验现象",
                            "body": "观察木块和铁块在水中的不同表现。",
                            "image_description": "浮力实验照片",
                            "image_url": f"/uploads/images/{filename}",
                        }
                    ],
                }
            )

            with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
                temp_file.write(pptx_bytes)
                temp_file.flush()

                from pptx import Presentation as PPTXPresentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                exported = PPTXPresentation(temp_file.name)
                self.assertEqual(len(exported.slides), 1)
                self.assertTrue(
                    any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in exported.slides[0].shapes)
                )
        finally:
            image_path.unlink(missing_ok=True)

    def test_title_subtitle_slide_does_not_embed_logo(self) -> None:
        UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"ppt-logo-{uuid.uuid4().hex}.png"
        image_path = UPLOADS_DIR / filename
        image_path.write_bytes(TINY_PNG)

        try:
            pptx_bytes = export_to_pptx(
                {
                    "title": "封面测试",
                    "classroom_script": "",
                    "slides": [
                        {
                            "template": "title_subtitle",
                            "title": "浮力探究",
                            "subtitle": "从生活现象进入课堂",
                        }
                    ],
                },
                {
                    "theme": "scholastic_blue",
                    "density": "balanced",
                    "school_name": "第一中学",
                    "logo_url": f"/uploads/images/{filename}",
                },
            )

            with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
                temp_file.write(pptx_bytes)
                temp_file.flush()

                from pptx import Presentation as PPTXPresentation
                from pptx.enum.shapes import MSO_SHAPE_TYPE

                exported = PPTXPresentation(temp_file.name)
                self.assertEqual(len(exported.slides), 1)
                self.assertFalse(
                    any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in exported.slides[0].shapes)
                )

                school_name_runs = [
                    run
                    for shape in exported.slides[0].shapes
                    if getattr(shape, "has_text_frame", False)
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.text == "第一中学"
                ]
                self.assertTrue(school_name_runs)
                self.assertGreaterEqual(school_name_runs[0].font.size.pt, 11)
        finally:
            image_path.unlink(missing_ok=True)

    def test_image_slide_does_not_overlay_placeholder_text_after_upload(self) -> None:
        UPLOADS_DIR.mkdir(parents=True, exist_ok=True)
        filename = f"ppt-image-{uuid.uuid4().hex}.png"
        image_path = UPLOADS_DIR / filename
        image_path.write_bytes(TINY_PNG)

        try:
            pptx_bytes = export_to_pptx(
                {
                    "title": "图片替换测试",
                    "classroom_script": "",
                    "slides": [
                        {
                            "template": "title_body_image",
                            "title": "实验照片",
                            "body": "先观察图片，再进入讲解。",
                            "image_description": "浮力实验装置",
                            "image_url": f"/uploads/images/{filename}",
                        }
                    ],
                }
            )

            with tempfile.NamedTemporaryFile(suffix=".pptx") as temp_file:
                temp_file.write(pptx_bytes)
                temp_file.flush()

                from pptx import Presentation as PPTXPresentation

                exported = PPTXPresentation(temp_file.name)
                slide_text = "\n".join(shape.text for shape in exported.slides[0].shapes if hasattr(shape, "text"))
                self.assertNotIn("浮力实验装置", slide_text)
                self.assertNotIn("图片占位", slide_text)
        finally:
            image_path.unlink(missing_ok=True)


if __name__ == "__main__":
    unittest.main()
